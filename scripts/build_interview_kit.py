"""Build the Week 4 interview kit (issue #24).

Outputs to output/interview_kit/: interviewer script + capture sheet
(.docx), a 6-slide mini deck (.pptx) — both upload straight into Google
Docs/Slides — and the three-report print pack. Regenerate any time the
measured numbers change.
"""

import shutil
from pathlib import Path

from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.util import Pt as PPt

OUT = Path("output/interview_kit")

# Measured results referenced in the materials (update if re-measured)
NUMBERS = {
    "recall": "9 of 10 planted discrepancies caught (7/7 deal-killers)",
    "false_positives": "0 false alarms on the clean company",
    "extraction": "100% clean scans / 100% office scans / 98.2% phone photos",
    "real_sweep": "12/12 real Companies House filings extracted, "
                  "every balance sheet ties",
    "cost": "roughly £8–10 of AI cost per data room",
}


def build_script_docx() -> Path:
    doc = Document()
    style = doc.styles["Normal"]
    style.font.name = "Calibri"
    style.font.size = Pt(11)

    doc.add_heading("Red Flag Report — Validation Interview Script", level=0)
    doc.add_paragraph(
        "Purpose: a price signal and a trust signal, from 5–10 people. "
        "This is NOT a sales meeting. Ask the two questions, capture answers "
        "verbatim, resist the urge to pitch or defend.")

    doc.add_heading("Who (recruit mix)", level=1)
    for line in (
            "3+ prospective or recent UK small-business buyers "
            "(UK ETA community, Rightbiz / BusinessesForSale enquirers, "
            "Searchfunder UK)",
            "2 accountants (ideally ones who do acquisition work — they are "
            "the factual audit)",
            "1–2 business-transfer brokers"):
        doc.add_paragraph(line, style="List Bullet")

    doc.add_heading("Setup (30 seconds, say roughly this)", level=1)
    doc.add_paragraph(
        '"I\'m testing a tool for people buying small UK businesses. You '
        "upload the seller's documents — accounts, VAT returns, bank "
        "statements, the lease — and it cross-checks them against each "
        "other and against Companies House, then produces this report. "
        'Every finding cites the exact page it came from. I\'d like your '
        'honest reaction to the report itself."')
    doc.add_paragraph(
        "Then hand over / share the MUTATED CAFÉ report first. Let them "
        "read. Do not narrate. Note what they linger on.")

    doc.add_heading("Demo order", level=1)
    for line in (
            "1. Mutated café report — the wow artifact (15 red flags: hidden "
            "loan, VAT mismatch, inflated P&L, month-9 lease break)",
            "2. Clean café report — the trust artifact (zero findings; point "
            "out the tool says so plainly rather than inventing problems)",
            "3. Chocolate Ice Cafe report — a real Companies House filing "
            "(distressed: £552k creditors vs £13k cash), showing honest "
            '"insufficient documents" behaviour on partial data'):
        doc.add_paragraph(line, style="List Bullet")

    doc.add_heading("THE TWO QUESTIONS (verbatim, in this order)", level=1)
    q = doc.add_paragraph()
    q.add_run('1.  "Would you have paid for this on your last deal — '
              'how much?"').bold = True
    doc.add_paragraph(
        "Allowed probes: silence (let them fill it); \"what would make that "
        "number higher?\" Do NOT suggest a price. If they ask what it costs, "
        'say "that\'s what I\'m trying to find out — what would it be worth '
        'to you?"')
    q2 = doc.add_paragraph()
    q2.add_run('2.  "What\'s missing that would make you distrust it?"'
               ).bold = True
    doc.add_paragraph(
        "Capture verbatim — these answers are the Month 2 backlog. For "
        "accountants add: \"is anything in this report factually wrong?\" "
        "and walk the findings with them.")

    doc.add_heading("Do NOT", level=1)
    for line in (
            "Pitch, explain the roadmap, or defend a criticism — say "
            '"that\'s really useful" and write it down',
            "Mention any price before they name one",
            "Show the technology / terminal — the report IS the product"):
        doc.add_paragraph(line, style="List Bullet")

    doc.add_heading("Evidence you can cite if asked how well it works",
                    level=1)
    for key in ("recall", "false_positives", "extraction", "real_sweep",
                "cost"):
        doc.add_paragraph(NUMBERS[key], style="List Bullet")

    doc.add_page_break()
    doc.add_heading("Capture sheet (one per interviewee)", level=0)
    table = doc.add_table(rows=8, cols=2)
    table.style = "Table Grid"
    fields = (
        "Name / role (buyer, accountant, broker)",
        "Background (last deal / sector / deal size)",
        "Q1 verbatim answer",
        "Q1 price named (£) — do not prompt",
        "Q2 verbatim answer (what would make them distrust it)",
        "Accountants only: factual errors found (exit criterion: zero)",
        "Unprompted reactions during reading (what they lingered on)",
        "Would they take a follow-up call / be a design partner? (Y/N)",
    )
    for row, label in zip(table.rows, fields, strict=True):
        row.cells[0].text = label
        row.cells[0].width = Pt(220)
    doc.add_paragraph()
    doc.add_paragraph(
        "Gate criteria (docs/04): GO needs ≥3 people independently at "
        "≥£300 and no accountant factual error. Converging on \"I'd only "
        "trust my accountant\" → pivot to accountant-tool model. Record "
        "answers in issue #19.")

    OUT.mkdir(parents=True, exist_ok=True)
    path = OUT / "interview_script_and_capture_sheet.docx"
    doc.save(str(path))
    return path


INK = RGBColor(0x1C, 0x1B, 0x1A)
RED = RGBColor(0xB3, 0x26, 0x1E)
GREY = RGBColor(0x6B, 0x66, 0x60)


def _slide(prs, title_text, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(11.5),
                                  Inches(1.1))
    p = tx.text_frame.paragraphs[0]
    p.text = title_text
    p.font.size = PPt(32)
    p.font.bold = True
    p.font.color.rgb = INK
    body = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(11.3),
                                    Inches(4.6))
    tf = body.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = text
        para.font.size = PPt(size)
        para.font.bold = bold
        para.font.color.rgb = color
        para.space_after = PPt(10)
    if note:
        notes = slide.notes_slide.notes_text_frame
        notes.text = note
    return slide


def build_deck_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _slide(prs, "The seller's documents disagree.\nFind out where.",
           [("Red Flag Report — automated due diligence for UK "
             "small-business purchases", 20, False, GREY),
            ("Every finding cited to the page it came from.", 18, False,
             GREY)],
           note="30 seconds max on this slide. The report does the talking.")

    _slide(prs, "Buying a UK company = inheriting its history",
           [("Most sub-£2m deals are share purchases — undisclosed debts, "
             "VAT arrears and payroll skeletons transfer to you.", 20, False,
             INK),
            ("The same revenue must appear four ways: VAT returns, filed "
             "accounts, management P&L, bank deposits.", 20, False, INK),
            ("Fraud and error live where those four disagree.", 22, True,
             RED)],
           note="This is the thesis: due diligence is reconciliation, "
                "not document reading.")

    _slide(prs, "What it does",
           [("1 · Reads the data room (accounts, VAT returns, bank "
             "statements, management P&L, lease)", 20, False, INK),
            ("2 · Cross-checks every figure — quarterly VAT triangle, "
             "Companies House charges register, filed vs management "
             "accounts, seller claims vs actuals", 20, False, INK),
            ("3 · Produces a cited Red Flag Report: what's wrong, what to "
             "ask the seller, what SPA warranty to demand", 20, False, INK),
            ("Unreadable figures are flagged for human review — never "
             "guessed.", 18, False, GREY)],
           note="Emphasise deterministic checks — the arithmetic is code, "
                "not AI opinion.")

    _slide(prs, "What it caught (test café, planted fraud)",
           [("£40k loan hidden from the accounts — found via Companies "
             "House charge + repayments still visible in the bank", 20,
             False, INK),
            ("Revenue inflated 11% above what was banked", 20, False, INK),
            ("VAT return understated vs the P&L for the same quarter", 20,
             False, INK),
            ("Lease terminable 9 months after completion", 20, False, INK),
            ("…and it stays silent on the honest twin: zero false alarms.",
             20, True, RED)],
           note="Hand out the two reports here — mutated first, then clean.")

    _slide(prs, "Measured, not promised",
           [(NUMBERS["recall"], 20, False, INK),
            (NUMBERS["false_positives"], 20, False, INK),
            ("Reading accuracy: " + NUMBERS["extraction"], 20, False, INK),
            (NUMBERS["real_sweep"], 20, False, INK)],
           note="If asked about method: synthetic data room with known "
                "planted discrepancies + real CH filings sweep.")

    _slide(prs, "Two questions",
           [("1 · Would you have paid for this on your last deal — "
             "how much?", 26, True, INK),
            ("2 · What's missing that would make you distrust it?", 26,
             True, INK)],
           note="Then stop talking. Capture verbatim.")

    path = OUT / "interview_mini_deck.pptx"
    prs.save(str(path))
    return path


def collect_report_pack() -> list[Path]:
    pack = OUT / "report_pack"
    pack.mkdir(parents=True, exist_ok=True)
    mapping = {
        "red_flag_report_copper_kettle_m1_clean.pdf":
            "1_test_cafe_with_planted_fraud.pdf",
        "red_flag_report_copper_kettle_clean_clean.pdf":
            "2_honest_test_cafe_zero_flags.pdf",
        "red_flag_report_real_chocolate_ice_clean.pdf":
            "3_real_companies_house_filing.pdf",
    }
    copied = []
    for src, dst in mapping.items():
        src_path = Path("output") / src
        if src_path.exists():
            shutil.copy(src_path, pack / dst)
            copied.append(pack / dst)
        else:
            print(f"  missing (regenerate with `diligence report`): {src}")
    return copied


if __name__ == "__main__":
    script = build_script_docx()
    deck = build_deck_pptx()
    pack = collect_report_pack()
    print(f"Interview kit in {OUT}/")
    print(f"  {script.name}")
    print(f"  {deck.name}")
    for p in pack:
        print(f"  report_pack/{p.name}")
